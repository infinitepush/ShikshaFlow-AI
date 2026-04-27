import os
import shutil
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont

OUTPUT_DIR = "output/"

def ensure_output_dir():
    """Create output folder if not exists."""
    if not os.path.exists(OUTPUT_DIR):
        os.makedirs(OUTPUT_DIR)
    return OUTPUT_DIR

def clear_output_dir():
    """Delete and recreate output folder (clean previous runs)."""
    if os.path.exists(OUTPUT_DIR):
        shutil.rmtree(OUTPUT_DIR)
    os.makedirs(OUTPUT_DIR)

def save_file(content, filename):
    """Save text or binary content to file."""
    ensure_output_dir()
    filepath = os.path.join(OUTPUT_DIR, filename)
    mode = 'wb' if isinstance(content, bytes) else 'w'
    with open(filepath, mode) as f:
        f.write(content)
    return filepath

def list_output_files():
    """List all generated files."""
    ensure_output_dir()
    return os.listdir(OUTPUT_DIR)

def _load_font(size, bold=False):
    font_candidates = [
        "arialbd.ttf" if bold else "arial.ttf",
        "DejaVuSans-Bold.ttf" if bold else "DejaVuSans.ttf",
    ]

    for font_name in font_candidates:
        try:
            return ImageFont.truetype(font_name, size)
        except OSError:
            continue

    return ImageFont.load_default()

def _draw_wrapped_text(draw, text, position, font, fill, max_width, line_spacing=12):
    x, y = position
    words = text.split()
    lines = []
    current_line = ""

    for word in words:
        test_line = f"{current_line} {word}".strip()
        bbox = draw.textbbox((0, 0), test_line, font=font)
        if bbox[2] - bbox[0] <= max_width:
            current_line = test_line
        else:
            if current_line:
                lines.append(current_line)
            current_line = word

    if current_line:
        lines.append(current_line)

    for line in lines:
        draw.text((x, y), line, fill=fill, font=font)
        bbox = draw.textbbox((x, y), line, font=font)
        y += (bbox[3] - bbox[1]) + line_spacing

    return y

def convert_pptx_to_images(pptx_path):
    """
    Converts slides in a PowerPoint file to images.
    This implementation renders basic slide content to images.
    """
    ensure_output_dir()  # Ensure output directory exists
    prs = Presentation(pptx_path)
    slide_images = []

    font_title = _load_font(58, bold=True)
    font_content = _load_font(36)

    for i, slide in enumerate(prs.slides):
        img_path = os.path.join(OUTPUT_DIR, f"slide_{i+1}.png")
        
        # Create an image with slide dimensions
        width, height = 1280, 720  # Standard HD resolution
        img = Image.new("RGB", (width, height), color=(255, 255, 255))  # White background
        draw = ImageDraw.Draw(img)
        
        # Try to add some basic text from the slide
        margin_x = 86
        y_offset = 76
        max_text_width = width - (margin_x * 2)
        
        # Add title if exists
        if slide.shapes.title and hasattr(slide.shapes.title, 'text'):
            title_text = slide.shapes.title.text.strip()
            if title_text:
                y_offset = _draw_wrapped_text(
                    draw,
                    title_text,
                    (margin_x, y_offset),
                    font_title,
                    (18, 24, 38),
                    max_text_width,
                    line_spacing=10,
                )
                y_offset += 46
        
        # Add content from placeholders
        for shape in slide.shapes:
            # Check if shape has text frame
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                # Get text content (ignore linter warning as this works at runtime)
                text_frame = shape.text_frame  # type: ignore
                if hasattr(text_frame, 'text'):
                    text = text_frame.text.strip()  # type: ignore
                    if text and shape != slide.shapes.title:  # Skip title as we already handled it
                        try:
                            lines = text.split('\n')
                            for line in lines:
                                if line.strip():
                                    clean_line = line.strip().lstrip("-").strip()
                                    clean_line = f"- {clean_line}"
                                    y_offset = _draw_wrapped_text(
                                        draw,
                                        clean_line,
                                        (margin_x + 12, y_offset),
                                        font_content,
                                        (45, 55, 72),
                                        max_text_width - 24,
                                        line_spacing=14,
                                    )
                                    y_offset += 20
                        except:
                            # Fallback if there are font issues
                            y_offset = _draw_wrapped_text(
                                draw,
                                text[:220],
                                (margin_x, y_offset),
                                font_content,
                                (45, 55, 72),
                                max_text_width,
                            )
        
        img.save(img_path)
        slide_images.append(img_path)
        print(f"Saved slide image: {img_path}, exists: {os.path.exists(img_path)}")

    # Upload images to Cloudinary
    try:
        from services.cloud_service import upload_file
        cloud_images = []
        for img_path in slide_images:
            upload_result = upload_file(img_path, resource_type="image")
            if upload_result:
                cloud_images.append(upload_result["secure_url"])
            else:
                # Fallback to local path if cloud upload fails
                cloud_images.append(img_path)
        # Return both local paths (for video creation) and cloud URLs (for response)
        return {"local_paths": slide_images, "cloud_urls": cloud_images}
    except Exception as e:
        print(f"Cloudinary upload failed: {e}")
        # Fallback to local paths if cloud upload fails
        return {"local_paths": slide_images, "cloud_urls": slide_images}
