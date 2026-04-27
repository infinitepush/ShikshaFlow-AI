import os

from pptx import Presentation

from utils.theme_utils import apply_theme


def generate_slides(slides_data, theme):
    output_dir = "output"
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    prs = Presentation()
    apply_theme(prs, theme)

    for slide in slides_data:
        slide_layout = prs.slide_layouts[1]
        s = prs.slides.add_slide(slide_layout)

        if s.shapes.title and "title" in slide:
            s.shapes.title.text = slide["title"]

        if len(s.placeholders) > 1:
            body = s.placeholders[1]
            if body.has_text_frame:
                text_frame = body.text_frame  # type: ignore
                text_frame.clear()

                for i, point in enumerate(slide["content"]):
                    bullet = f"- {point}"
                    if i == 0:
                        text_frame.text = bullet
                    else:
                        p = text_frame.add_paragraph()
                        p.text = bullet

    path = "output/slides.pptx"
    prs.save(path)
    return path
